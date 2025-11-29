#!/usr/bin/env python3
"""
Script to create PowerPoint presentation from the microservices architecture content.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Microservices TODO Application"
    subtitle.text = "Containerized, Automated & Secure Deployment\n\nComplete DevOps Pipeline Implementation\nTerraform + Ansible + Docker + CI/CD + Drift Detection"
    
    # Slide 2: Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Architecture Overview"
    content.text = """6 Microservices in Different Languages:
• Frontend (Vue.js)
• Auth API (Go) 
• Todos API (Node.js)
• Users API (Java Spring Boot)
• Log Processor (Python)
• Redis Queue

Key Features:
• Traefik Reverse Proxy with SSL termination
• Docker Containerization for all services
• Redis for message queuing"""

    # Slide 3: PART 1 - Containerization
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "PART 1: Application Containerization"
    content.text = """Implementation Flow:
• Individual Dockerfiles for each service
• Docker Compose orchestrates all services  
• Traefik Proxy handles SSL + Routing

Single Command Deployment:
docker compose up -d

HTTPS Endpoints:
• https://your-domain.com
• https://your-domain.com/api/auth
• https://your-domain.com/api/todos
• https://your-domain.com/api/users

Features:
• Automatic SSL with Let's Encrypt certificates
• HTTP → HTTPS automatic redirection"""

    # Slide 4: PART 2 - Infrastructure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "PART 2: Infrastructure & Automation"
    content.text = """Three Main Components:

Terraform:
• AWS EC2 provisioning
• Security groups configuration
• Remote state management (S3)
• Dynamic Ansible inventory generation

Ansible:
• Dependencies role (Docker, Git, packages)
• Deploy role (app deployment, SSL setup)
• Idempotent operations

CI/CD:
• GitHub Actions workflows
• Drift detection with email alerts
• Automated deployment pipelines"""

    # Slide 5: Terraform Implementation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Terraform: Idempotent Infrastructure"
    content.text = """Key Components:
• main.tf: EC2 instances + Security Groups
• Remote State: S3 Backend for team collaboration
• Dynamic Inventory: Auto-generated for Ansible

Automatic Ansible Execution:
resource "null_resource" "run_ansible" {
  provisioner "local-exec" {
    command = "ansible-playbook -i inventory.ini site.yml"
  }
}

Features:
• Fully idempotent operations
• No resource recreation unless drift occurs
• Automatic integration with Ansible"""

    # Slide 6: Ansible Roles
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Ansible: Server Configuration & Deployment"
    content.text = """Dependencies Role:
• Docker & Docker Compose installation
• Git and system packages
• User permissions and groups

Deploy Role:
• Repository cloning and updates
• Container management with Docker Compose
• SSL configuration and health checks

Key Features:
• Idempotent Deployment: No restart unless files changed
• Git Integration: Automatic repo cloning and updates
• Container Orchestration: Full Docker Compose management"""

    # Slide 7: CI/CD Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CI/CD: Automated Deployment Pipeline"
    content.text = """Two Workflow Types:

Infrastructure Workflow:
• Triggers on infra/terraform/** changes
• Triggers on infra/ansible/** changes

Application Workflow:
• Triggers on service code changes
• Triggers on docker-compose.yml changes

Key Features:
• Drift Detection: terraform plan → email alert → manual approval
• Conditional Deployment: Only runs if changes detected
• Email Notifications: Automatic alerts for infrastructure drift"""

    # Slide 8: Drift Detection
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Drift Detection & Safety Mechanism"
    content.text = """5-Step Safety Process:

1. Plan Check: terraform plan -detailed-exitcode
2. Drift Detection: Exit code = 2 indicates changes
3. Email Alert: Notify stakeholders, pause pipeline
4. Manual Approval: GitHub Environment protection
5. Apply Changes: terraform apply only after approval

Safety Rules:
• No Drift = Automatic Proceed
• Drift Detected = Email + Manual Approval Required

This ensures complete transparency and control over infrastructure changes."""

    # Slide 9: Security Implementation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Security Implementation"
    content.text = """Multi-Layer Security Approach:

Network Security:
• AWS Security Groups with minimal port exposure
• SSH key-based authentication
• Only ports 22, 80, 443 exposed

SSL/TLS Security:
• Let's Encrypt certificates with automatic renewal
• HTTP → HTTPS automatic redirection
• Traefik handles SSL termination

Application Security:
• JWT token authentication
• API authorization middleware
• Container isolation and network segmentation"""

    # Slide 10: Single Command Deployment
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "PART 3: Single Command Deployment"
    content.text = """Complete Stack Deployment:
terraform apply -auto-approve

5-Step Automated Process:
1. Provision: AWS EC2 + Security Groups
2. Generate: Ansible inventory file
3. Configure: Install all dependencies
4. Deploy: Application containers
5. Secure: Traefik + SSL setup

Benefits:
• Fully Automated: Zero manual intervention required
• Idempotent: Skip unchanged resources
• Production Ready: HTTPS endpoints immediately available"""

    # Slide 11: Key Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Benefits & Achievements"
    content.text = """🚀 Automation:
• Single command deployment
• Zero manual configuration
• Automatic SSL setup

🛡️ Safety:
• Infrastructure drift detection
• Email notifications for changes
• Manual approval gates

🔄 Reliability:
• Idempotent operations
• Remote state management
• Rollback capability

📈 Scalability:
• Microservices architecture
• Container orchestration
• Cloud-native design"""

    # Slide 12: Demo & Questions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Live Demonstration & Questions"
    content.text = """Ready to Deploy!

Commands to Execute:
cd infra/terraform
terraform apply -auto-approve

Application Access:
• https://your-domain.com
• https://your-domain.com/api/auth
• https://your-domain.com/api/todos
• https://your-domain.com/api/users

Questions & Discussion
Thank you for your attention!"""

    # Save presentation
    prs.save('microservice-arch.pptx')
    print("PowerPoint presentation created: microservice-arch.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: python-pptx library not installed.")
        print("Install it with: pip install python-pptx")
        print("Then run this script again.")
    except Exception as e:
        print(f"Error creating presentation: {e}")
        print("Please check the microservice-arch-outline.txt file for manual PowerPoint creation.")